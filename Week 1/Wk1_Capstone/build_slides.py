#!/usr/bin/env python
"""
build_slides.py — generates group3.pptx from the executed capstone notebook.

Usage (from repo root):
    .venv/bin/python "Week 1/Wk1_Capstone/build_slides.py"

Edit DATA dict to update numbers after each full-data run.
Change CONTRIB_SEED to re-roll the contribution matrix.
"""

import base64, json, pathlib, sys
from collections import Counter

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Paths ─────────────────────────────────────────────────────────────────────
HERE       = pathlib.Path(__file__).parent
ASSETS_DIR = HERE / "slides_assets"
NB_PATH    = HERE / "capstone_churn_executed-fullset.ipynb"
OUT_PATH   = HERE / "group3.pptx"
ASSETS_DIR.mkdir(exist_ok=True)

# ── DATA — update values after full-data run ──────────────────────────────────
DATA = {
    "group":        "Group 3",
    "course":       "ICT6001C — AI Programme | Week 1 Capstone",
    "repo_link":    "https://github.com/kinwahlai/ICT6001C",
    "youtube_link": "<YouTube link>",
    "members": [
        {"name": "YEO XIU JUAN",   "id": "2570472"},
        {"name": "LIM KHOON SENG", "id": "2570459"},
        {"name": "LAI KIN WAH",    "id": "2570461"},
    ],

    # ── CV results (5-fold, defaults, full 594k data) ─────────────────────────
    "cv_results": [
        {"model": "Logistic Regression", "roc_auc": "0.9059±0.0011", "recall": "0.8808±0.0027"},
        {"model": "Random Forest",       "roc_auc": "0.8872±0.0013", "recall": "0.7072±0.0024"},
        {"model": "XGBoost",             "roc_auc": "0.9132±0.0012", "recall": "0.8245±0.0038"},
    ],

    # ── Champion ──────────────────────────────────────────────────────────────
    "champion":          "XGBoost",
    "champion_family":   "Boosting",
    "champion_params":   "learning_rate=0.05, max_depth=3, n_estimators=400",
    "champion_cv_auc":   "0.9104±0.0014",
    "champion_cv_recall":"0.8582±0.0032",

    # ── Hold-out (evaluated once, full 594k data) ─────────────────────────────
    "threshold":    "0.61",
    "hold_auc":     "0.9132",
    "hold_recall":  "0.8079",
    "hold_bacc":    "0.8292",

    # ── EDA insights ─────────────────────────────────────────────────────────
    "eda_insights": [
        "Tenure: New customers (<12 months) highest churn risk — focus retention on year 1.",
        "MonthlyCharges: Churners pay ~$15–20/month more (pricing sensitivity).",
        "Contract type: Month-to-month churn ~42%; two-year contracts <3%.",
    ],

    # ── Business narrative ────────────────────────────────────────────────────
    "biz_reco": (
        "Tier retention offers by predicted churn probability — "
        "concentrate premium spend on highest-risk customers."
    ),
}

MEMBERS = [m["name"] for m in DATA["members"]]

# ── Fixed contribution matrix ─────────────────────────────────────────────────
CONTRIBUTIONS = [
    {"area": "Part A: Pipeline Engineering",  "primary": "LIM KHOON SENG", "secondary": "YEO XIU JUAN"},
    {"area": "Part A: EDA & Data Quality",    "primary": "YEO XIU JUAN",   "secondary": "LAI KIN WAH"},
    {"area": "Part B: Model Selection",       "primary": "YEO XIU JUAN",   "secondary": "LIM KHOON SENG"},
    {"area": "Part C: Hyperparameter Tuning", "primary": "LIM KHOON SENG", "secondary": "YEO XIU JUAN"},
    {"area": "Part E: Business Decision",     "primary": "YEO XIU JUAN",   "secondary": "LAI KIN WAH"},
    {"area": "Report & Slides",               "primary": "LAI KIN WAH",    "secondary": "LIM KHOON SENG"},
]

# ── Theme ─────────────────────────────────────────────────────────────────────
W      = Inches(13.333)   # 16:9 widescreen width
H      = Inches(7.5)
BAR_H  = Inches(0.72)
MARGIN = Inches(0.45)

C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_ACCENT = RGBColor(0x29, 0x80, 0xB9)   # notebook blue
C_TITLE  = RGBColor(0x1A, 0x56, 0x76)   # darker blue for title slide bar
C_TEXT   = RGBColor(0x2C, 0x3E, 0x50)   # dark slate
C_LIGHT  = RGBColor(0xEB, 0xF5, 0xFB)   # alternating table row tint
C_GREEN  = RGBColor(0x1E, 0x8B, 0x4C)
C_RED    = RGBColor(0xC0, 0x39, 0x2B)
C_GRAY   = RGBColor(0x7F, 0x8C, 0x8D)

FONT = "Calibri"

# ── Helpers ───────────────────────────────────────────────────────────────────
def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])   # Blank layout

def set_bg(sl, color=C_WHITE):
    fill = sl.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(sl, l, t, w, h, fill, line=None):
    sh = sl.shapes.add_shape(1, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line: sh.line.color.rgb = line
    else:    sh.line.fill.background()
    return sh

def txbox(sl, text, l, t, w, h, size=14, bold=False,
          color=C_TEXT, align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = sl.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = wrap
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.name = FONT; r.font.size = Pt(size)
    r.font.bold = bold; r.font.color.rgb = color; r.font.italic = italic
    return tb

def bullets(sl, items, l, t, w, h, size=14, color=C_TEXT, marker="◆"):
    tb = sl.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(3)
        r = p.add_run()
        r.text = f"{marker}  {item}"
        r.font.name = FONT; r.font.size = Pt(size); r.font.color.rgb = color
    return tb

def title_bar(sl, title, sub=None):
    """Blue bar at top; returns y-offset for content start."""
    rect(sl, 0, 0, W, BAR_H, C_ACCENT)
    txbox(sl, title, MARGIN, Inches(0.1), W - MARGIN*2, BAR_H - Inches(0.15),
          size=22, bold=True, color=C_WHITE)
    y = BAR_H
    if sub:
        txbox(sl, sub, MARGIN, y, W - MARGIN*2, Inches(0.3),
              size=10, color=C_GRAY, italic=True)
        y += Inches(0.32)
    return y + Inches(0.1)

def table(sl, headers, rows, l, t, w, col_widths=None, hdr_size=11, row_size=10):
    row_h = Inches(0.34)
    tot_h = row_h * (len(rows) + 1)
    gf    = sl.shapes.add_table(len(rows)+1, len(headers), l, t, w, tot_h)
    tbl   = gf.table

    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = cw

    for ci, hdr in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.text = hdr
        cell.fill.solid(); cell.fill.fore_color.rgb = C_ACCENT
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.runs[0] if p.runs else p.add_run()
        run.font.name = FONT; run.font.size = Pt(hdr_size)
        run.font.bold = True; run.font.color.rgb = C_WHITE

    for ri, row in enumerate(rows):
        bg = C_LIGHT if ri % 2 == 0 else C_WHITE
        for ci, val in enumerate(row):
            cell = tbl.cell(ri+1, ci)
            cell.text = str(val)
            cell.fill.solid(); cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            run = p.runs[0] if p.runs else p.add_run()
            run.font.name = FONT; run.font.size = Pt(row_size)
            run.font.color.rgb = C_TEXT

    return tbl

def picture(sl, path, l, t, w=None, h=None):
    if w and h:   sl.shapes.add_picture(str(path), l, t, width=w, height=h)
    elif w:       sl.shapes.add_picture(str(path), l, t, width=w)
    elif h:       sl.shapes.add_picture(str(path), l, t, height=h)

# ── Image extraction ──────────────────────────────────────────────────────────
def extract_images(nb_path):
    nb   = json.loads(pathlib.Path(nb_path).read_text())
    imgs = {}
    idx  = 0
    for cell in nb["cells"]:
        for out in cell.get("outputs", []):
            raw = out.get("data", {}).get("image/png")
            if raw is None:
                continue
            if isinstance(raw, list):
                raw = "".join(raw)
            path = ASSETS_DIR / f"fig{idx}.png"
            path.write_bytes(base64.b64decode(raw))
            imgs[idx] = {"cell": cell["id"], "path": path}
            print(f"  fig{idx}.png  ← cell {cell['id']}")
            idx += 1
    return imgs
#   Expected order: 0=churn-dist, 1=key-drivers, 2=CV-bar, 3=recall-vs-threshold, 4=confusion-matrix

# ── Slide builders ────────────────────────────────────────────────────────────
def s1_title(prs):
    sl = blank(prs); set_bg(sl)
    rect(sl, 0, 0, W, Inches(2.5), C_TITLE)
    txbox(sl, "Telco Customer Churn Prediction",
          MARGIN, Inches(0.2), W - MARGIN*2, Inches(1.1),
          size=36, bold=True, color=C_WHITE)
    txbox(sl, DATA["course"],
          MARGIN, Inches(1.35), W - MARGIN*2, Inches(0.45),
          size=16, color=C_WHITE)
    txbox(sl, DATA["group"],
          MARGIN, Inches(1.85), W - MARGIN*2, Inches(0.4),
          size=14, bold=True, color=RGBColor(0xAE, 0xD6, 0xF1))

    y = Inches(2.85)
    for m in DATA["members"]:
        txbox(sl, f"  {m['name']} (ID: {m['id']})",
              MARGIN, y, Inches(6), Inches(0.38), size=15, color=C_TEXT)
        y += Inches(0.4)

    txbox(sl, f"Repo :  {DATA['repo_link']}",
          MARGIN, Inches(6.4), W - MARGIN*2, Inches(0.35), size=11, color=C_ACCENT)
    txbox(sl, f"Video:  {DATA['youtube_link']}",
          MARGIN, Inches(6.8), W - MARGIN*2, Inches(0.35), size=11, color=C_ACCENT)


def s2_exec_summary(prs):
    d  = DATA
    sl = blank(prs); set_bg(sl)
    y0 = title_bar(sl, "Executive Summary")
    b  = [
        f"Dataset: 594,194 Telco customers  |  Churn rate ~22.5%  |  Class imbalance handled with SMOTE",
        f"3 model families compared via 5-fold stratified CV (ROC-AUC primary, Recall secondary)",
        f"Champion: {d['champion']} ({d['champion_family']}) — highest CV ROC-AUC",
        f"Champion tuned via GridSearchCV (champion-only, ≤50-config budget); CV ROC-AUC {d['champion_cv_auc']}",
        f"Hold-out (evaluated once):  ROC-AUC {d['hold_auc']}  ·  Recall {d['hold_recall']}  at threshold {d['threshold']}",
        f"Threshold {d['threshold']} = highest cut still holding ≥80% recall; FN far costlier than FP, and SMOTE lifts probabilities above 0.50",
        d["biz_reco"],
    ]
    bullets(sl, b, MARGIN, y0, W - MARGIN*2, H - y0 - MARGIN, size=15)


def s3_eda(prs, imgs):
    sl = blank(prs); set_bg(sl)
    y0 = title_bar(sl, "Data Engineering & EDA  (Part A)",
                   "CRISP-DM: Data Preparation  |  80/20 stratified split  |  X_test locked until §4")

    left_w = Inches(4.8)
    txbox(sl, "EDA Insights", MARGIN, y0, left_w, Inches(0.35),
          size=13, bold=True, color=C_ACCENT)
    extra = [
        "No missing values; TotalCharges coerced from blank strings.",
        "Churn rate ~22.5% — moderate imbalance addressed by SMOTE inside pipeline.",
        "80/20 stratified train/test split; X_test untouched until §4.",
    ]
    bullets(sl, DATA["eda_insights"] + extra,
            MARGIN, y0 + Inches(0.4), left_w, H - y0 - Inches(0.6), size=12)

    if 1 in imgs:
        chart_x = left_w + MARGIN * 2
        picture(sl, imgs[1]["path"], chart_x, y0,
                w=W - chart_x - MARGIN)


def s4_pipeline(prs):
    sl = blank(prs); set_bg(sl)
    y0 = title_bar(sl, "Leakage-Safe Pipeline Architecture  (Part A)",
                   "CRISP-DM: Data Preparation  |  All transforms encapsulated inside sklearn Pipeline")

    # ── Visual box-flow: 3 stage boxes + arrows ──────────────────────────────
    box_t  = y0
    box_h  = Inches(2.2)
    arr_w  = Inches(0.45)
    total  = W - MARGIN * 2
    # box widths: narrow | narrow | wide
    bw1 = Inches(2.8); bw2 = Inches(2.2); bw3 = total - bw1 - bw2 - arr_w * 2
    x1 = MARGIN; x2 = x1 + bw1 + arr_w; x3 = x2 + bw2 + arr_w

    # Stage 1 — BEFORE SPLIT
    rect(sl, x1, box_t, bw1, Inches(0.38), C_ACCENT)
    txbox(sl, "BEFORE SPLIT", x1 + Inches(0.08), box_t + Inches(0.06),
          bw1 - Inches(0.1), Inches(0.3), size=11, bold=True, color=C_WHITE)
    rect(sl, x1, box_t + Inches(0.38), bw1, box_h - Inches(0.38), C_LIGHT)
    bullets(sl, ["Coerce TotalCharges to numeric", "Map target: Yes → 1, No → 0",
                 "Drop ID column"],
            x1 + Inches(0.08), box_t + Inches(0.48), bw1 - Inches(0.12),
            box_h - Inches(0.55), size=11, marker="•")

    # Arrow 1
    txbox(sl, "→", x1 + bw1, box_t + box_h * 0.4, arr_w, Inches(0.4),
          size=20, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

    # Stage 2 — SPLIT
    rect(sl, x2, box_t, bw2, Inches(0.38), C_ACCENT)
    txbox(sl, "80 / 20 SPLIT", x2 + Inches(0.08), box_t + Inches(0.06),
          bw2 - Inches(0.1), Inches(0.3), size=11, bold=True, color=C_WHITE)
    rect(sl, x2, box_t + Inches(0.38), bw2, box_h - Inches(0.38), C_LIGHT)
    bullets(sl, ["Stratified by Churn", "Seed = 42",
                 "Test set locked until §4"],
            x2 + Inches(0.08), box_t + Inches(0.48), bw2 - Inches(0.12),
            box_h - Inches(0.55), size=11, marker="•")

    # Arrow 2
    txbox(sl, "→", x2 + bw2, box_t + box_h * 0.4, arr_w, Inches(0.4),
          size=20, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

    # Stage 3 — INSIDE PIPELINE (per CV fold)
    rect(sl, x3, box_t, bw3, Inches(0.38), C_TITLE)
    txbox(sl, "INSIDE PIPELINE  (per CV fold)", x3 + Inches(0.08), box_t + Inches(0.06),
          bw3 - Inches(0.1), Inches(0.3), size=11, bold=True, color=C_WHITE)
    rect(sl, x3, box_t + Inches(0.38), bw3, box_h - Inches(0.38), C_LIGHT)
    pipe_items = [
        "Numeric  →  StandardScaler",
        "Categorical  →  OneHotEncoder  →  χ² top-50%",
        "↓  SMOTE  (training fold only)",
        "↓  XGBoost champion",
    ]
    bullets(sl, pipe_items, x3 + Inches(0.08), box_t + Inches(0.48),
            bw3 - Inches(0.12), box_h - Inches(0.55), size=12, marker="")

    # Design Decisions (below the flow)
    y_note = box_t + box_h + Inches(0.25)
    points = [
        "SMOTE runs inside the pipeline — oversampling touches training folds only, never validation/test.",
        "χ² feature selection keeps the top-50% of encoded categories most associated with churn.",
        "Numeric features scaled so Logistic Regression converges; trees are scale-invariant.",
        "Every transform is fit on the training fold only — the test set is untouched until §4.",
    ]
    txbox(sl, "Design Decisions", MARGIN, y_note, W - MARGIN*2, Inches(0.35),
          size=13, bold=True, color=C_ACCENT)
    bullets(sl, points, MARGIN, y_note + Inches(0.38), W - MARGIN*2,
            H - y_note - Inches(0.5), size=12)


def s5_model_selection(prs, imgs):
    sl = blank(prs); set_bg(sl)
    y0 = title_bar(sl, "Champion Model Selection  (Part B)",
                   "5-fold stratified CV at representative defaults  |  ROC-AUC primary, Recall secondary")

    # CV results table (left half)
    headers = ["Model", "ROC-AUC (mean±std)", "Recall (mean±std)"]
    rows    = [(r["model"], r["roc_auc"], r["recall"]) for r in DATA["cv_results"]]
    col_w   = [Inches(3.5), Inches(2.3), Inches(2.3)]
    tbl_w   = sum(col_w)
    table(sl, headers, rows, MARGIN, y0, tbl_w, col_widths=col_w)

    # Metric rationale (below table)
    y_rat = y0 + Inches(1.6)
    pts = [
        "ROC-AUC: threshold-independent ranking quality — fair comparison across all families.",
        "Recall: tracks business cost of missed churners (FN). F1/Precision omitted — threshold is set in §4.",
        "SMOTE oversampling inside ImbPipeline — minority class boosted only in training folds.",
    ]
    bullets(sl, pts, MARGIN, y_rat, tbl_w + Inches(0.1), H - y_rat - MARGIN, size=12)

    # CV bar chart (right)
    if 2 in imgs:
        chart_x = tbl_w + MARGIN * 2
        picture(sl, imgs[2]["path"], chart_x, y0,
                w=W - chart_x - MARGIN)


def s6_champion(prs):
    d  = DATA
    sl = blank(prs); set_bg(sl)
    y0 = title_bar(sl, f"Champion Justification: {d['champion']}  (Part B)",
                   "Selected by highest 5-fold CV ROC-AUC at default params")

    # Champion summary block
    rect(sl, MARGIN, y0, W - MARGIN*2, Inches(1.3), C_LIGHT)
    champ_txt = (
        f"Champion Model:  {d['champion']}  ({d['champion_family']} family)\n"
        f"CV ROC-AUC:      {d['champion_cv_auc']}      "
        f"CV Recall:  {d['champion_cv_recall']}"
    )
    tb = sl.shapes.add_textbox(MARGIN + Inches(0.15), y0 + Inches(0.1),
                               W - MARGIN*2 - Inches(0.3), Inches(1.1))
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(champ_txt.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(6)
        r = p.add_run(); r.text = line
        r.font.name = FONT; r.font.size = Pt(15)
        r.font.bold = (i == 0); r.font.color.rgb = C_TEXT

    # Justification bullets
    y_j = y0 + Inches(1.5)
    txbox(sl, "Why this champion?", MARGIN, y_j, W - MARGIN*2, Inches(0.35),
          size=13, bold=True, color=C_ACCENT)
    pts = [
        f"{d['champion']} achieved the highest mean CV ROC-AUC across all 3 families.",
        "Low standard deviation across folds → consistent performance, low overfitting risk.",
        "ROC-AUC is threshold-independent — a fair selection criterion before §4 threshold shift.",
        f"{d['champion_family']} family: sequentially corrects errors — strongest on tabular, imbalanced data.",
        "Provides strong baseline for GridSearchCV tuning in §3 (champion-only grid search).",
    ]
    bullets(sl, pts, MARGIN, y_j + Inches(0.4), W - MARGIN*2, Inches(2.1), size=14)

    # 3-family comparison table (stacked below bullets, full width)
    y_tbl = y_j + Inches(2.6)
    txbox(sl, "Model families considered", MARGIN, y_tbl, W - MARGIN*2, Inches(0.35),
          size=13, bold=True, color=C_ACCENT)
    headers = ["Family", "Model", "Rationale"]
    rows    = [
        ("Linear",  "Logistic Regression", "Interpretable baseline; linear decision boundary"),
        ("Bagging", "Random Forest",       "Variance reduction; built-in feature importance"),
        ("Boosting","XGBoost",             "Sequentially corrects errors; strong on tabular data"),
    ]
    col_w = [Inches(1.8), Inches(2.8), Inches(7.8)]
    tbl_w = sum(col_w)
    table(sl, headers, rows, MARGIN, y_tbl + Inches(0.4), tbl_w, col_widths=col_w)


def s7_tuning(prs):
    d  = DATA
    sl = blank(prs); set_bg(sl)
    y0 = title_bar(sl, "Hyperparameter Tuning  (Part C)",
                   "GridSearchCV on champion only  |  ≤50 candidate-config cap")

    # Strategy table (left column)
    left_w = Inches(7.5)
    headers = ["Decision", "Choice", "Rationale"]
    rows    = [
        ("Search method",   "GridSearchCV",          "Exhaustive over small grids; reproducible"),
        ("Scope",           "Champion only",         "One champion declared in §2; tune alone"),
        ("Iterations cap",  "≤50 candidate configs", f"Champion grid ≤12 configs (well within cap)"),
        ("CV folds",        "5-fold stratified",     "Validation repeats, not search iterations"),
        ("Tune dataset",    "100k subsample",        "4–5× faster; refit on full 475k rows"),
        ("refit metric",    "roc_auc",               "Threshold-independent; recall tuned in §4"),
    ]
    col_w = [Inches(2.0), Inches(2.2), Inches(3.3)]
    table(sl, headers, rows, MARGIN, y0, left_w, col_widths=col_w)

    # Tuning log (right column)
    log_x = left_w + MARGIN * 2
    log_w = W - log_x - MARGIN
    txbox(sl, "Tuning Log", log_x, y0, log_w, Inches(0.35),
          size=13, bold=True, color=C_ACCENT)
    log_lines = [
        f"Champion:     {d['champion']}",
        f"Best params:  {d['champion_params']}",
        f"CV ROC-AUC:  {d['champion_cv_auc']}",
        f"CV Recall:    {d['champion_cv_recall']}",
        "",
        "Refit best config on full 475k-row train set",
    ]
    tb = sl.shapes.add_textbox(log_x, y0 + Inches(0.4), log_w, Inches(3.0))
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(log_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(3)
        r = p.add_run(); r.text = line
        r.font.name = FONT; r.font.size = Pt(12); r.font.color.rgb = C_TEXT


def s8_stability(prs, imgs):
    d  = DATA
    sl = blank(prs); set_bg(sl)
    y0 = title_bar(sl, "Final Model Stability  (Part B → Part E)",
                   "CV cross-validated scores (train)  +  Hold-out evaluation (test, accessed once)")

    left_w = Inches(6.0)
    # CV scores table
    txbox(sl, "5-fold CV Scores (defaults)", MARGIN, y0, left_w, Inches(0.35),
          size=13, bold=True, color=C_ACCENT)
    cv_headers = ["Model", "ROC-AUC", "Recall"]
    cv_rows    = [(r["model"], r["roc_auc"], r["recall"]) for r in d["cv_results"]]
    col_w      = [Inches(3.2), Inches(1.5), Inches(1.3)]
    table(sl, cv_headers, cv_rows, MARGIN, y0 + Inches(0.38), left_w, col_widths=col_w)

    # Hold-out metrics
    y_ho = y0 + Inches(2.0)
    txbox(sl, "Hold-Out Test Metrics  (threshold = " + d['threshold'] + ")",
          MARGIN, y_ho, left_w, Inches(0.35), size=13, bold=True, color=C_ACCENT)
    ho_headers = ["Metric", "Value"]
    ho_rows    = [
        ("ROC-AUC",           d["hold_auc"]),
        ("Recall",            d["hold_recall"]),
        ("Balanced Accuracy", d["hold_bacc"]),
    ]
    table(sl, ho_headers, ho_rows, MARGIN, y_ho + Inches(0.38), Inches(3.5),
          col_widths=[Inches(2.3), Inches(1.2)])

    # Confusion matrix image (right)
    if 4 in imgs:
        chart_x = left_w + MARGIN * 2
        picture(sl, imgs[4]["path"], chart_x, y0,
                w=W - chart_x - MARGIN)


def s9_business(prs, imgs):
    d  = DATA
    sl = blank(prs); set_bg(sl)
    y0 = title_bar(sl, "Business Decision Reasoning  (Part E)",
                   "CRISP-DM: Deployment  |  Asymmetric error costs  |  Threshold shift")

    left_w = Inches(6.2)
    # Cost table
    txbox(sl, "Asymmetric Error Costs", MARGIN, y0, left_w, Inches(0.35),
          size=13, bold=True, color=C_ACCENT)
    cost_h = ["Error Type", "Consequence", "Severity"]
    cost_r = [
        ("False Negative — miss a churner",
         "Customer leaves uncontacted → lost lifetime value", "CRITICAL"),
        ("False Positive — flag loyal customer",
         "Wasted retention offer — comparatively cheap, recoverable", "Moderate"),
    ]
    table(sl, cost_h, cost_r, MARGIN, y0 + Inches(0.38), left_w,
          col_widths=[Inches(2.5), Inches(2.8), Inches(0.9)])

    # Decision points
    y_dec = y0 + Inches(1.75)
    dec = [
        "FN far costlier than FP → prioritise recall; require recall ≥ 80%.",
        f"Among thresholds meeting the floor, {d['threshold']} is the highest — limits false positives.",
        "SMOTE inflates predicted probabilities, so the operating point sits above 0.50, not below.",
        f"At threshold {d['threshold']}: Recall {d['hold_recall']} on hold-out — "
        f"{float(d['hold_recall'])*100:.0f}% of real churners flagged for intervention.",
        d["biz_reco"],
    ]
    txbox(sl, "Threshold Decision", MARGIN, y_dec, left_w, Inches(0.35),
          size=13, bold=True, color=C_ACCENT)
    bullets(sl, dec, MARGIN, y_dec + Inches(0.38), left_w,
            H - y_dec - Inches(0.5), size=12)

    # Recall-vs-Threshold chart (right)
    if 3 in imgs:
        chart_x = left_w + MARGIN * 2
        picture(sl, imgs[3]["path"], chart_x, y0,
                w=W - chart_x - MARGIN)


def s10_contributions(prs):
    sl = blank(prs); set_bg(sl)
    y0 = title_bar(sl, "Individual Contributions")

    # Name / ID reference table
    id_headers = ["Name", "Student ID"]
    id_rows    = [(m["name"], m["id"]) for m in DATA["members"]]
    table(sl, id_headers, id_rows, MARGIN, y0, Inches(4.5),
          col_widths=[Inches(3.0), Inches(1.5)])

    # Contributions table (names only, no repeated IDs)
    y_contrib = y0 + Inches(1.6)
    headers = ["Work Area", "Primary Owner", "Secondary Support"]
    rows    = [(c["area"], c["primary"], c["secondary"]) for c in CONTRIBUTIONS]
    col_w   = [Inches(4.5), Inches(3.5), Inches(3.5)]
    tbl_w   = sum(col_w)
    table(sl, headers, rows, MARGIN, y_contrib, tbl_w, col_widths=col_w)



# ── Main ──────────────────────────────────────────────────────────────────────
def main():
    if not NB_PATH.exists():
        sys.exit(f"ERROR: {NB_PATH} not found. Run: .venv/bin/python run_nb.py capstone_churn.ipynb")

    print("Extracting images from executed notebook…")
    imgs = extract_images(NB_PATH)
    print(f"  → {len(imgs)} image(s) extracted to {ASSETS_DIR}/\n")
    if len(imgs) < 5:
        print(f"  WARNING: expected 5 images, got {len(imgs)}. "
              "Some slides may be missing charts.")

    print("Assembling slides…")
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    s1_title(prs)
    s2_exec_summary(prs)
    s3_eda(prs, imgs)
    s4_pipeline(prs)
    s5_model_selection(prs, imgs)
    s6_champion(prs)
    s7_tuning(prs)
    s8_stability(prs, imgs)
    s9_business(prs, imgs)
    s10_contributions(prs)

    prs.save(OUT_PATH)
    print(f"\nSaved: {OUT_PATH}  ({len(prs.slides)} slides)")
    print("\nContribution matrix:")
    for c in CONTRIBUTIONS:
        print(f"  {c['area']:<35}  Primary: {c['primary']:<15}  Secondary: {c['secondary']}")

    print("\nNext steps:")
    print("  1. Open group3.pptx in PowerPoint / Keynote to review & edit.")
    print("  2. Fill <repo link> and <YouTube link> in DATA at top of this script, re-run.")
    print("  3. Export group3.pdf for submission (filename must be groupX.pdf).")
    print("  4. Re-run this script after the full 594k notebook run to get final numbers.")


if __name__ == "__main__":
    main()
